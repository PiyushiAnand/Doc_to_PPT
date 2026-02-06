import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ---------------------------------------------------------
# CONFIGURATION: MODERN GEOMETRIC THEME
# ---------------------------------------------------------
# 1. Color Palette
PRIMARY_COLOR = RGBColor(46, 0, 75)     # Dark Indigo/Violet (Covers/Headers)
ACCENT_COLOR = RGBColor(255, 100, 80)   # Pink-to-Orange proxy (Coral) for Lines
ICON_COLOR = RGBColor(0, 180, 216)      # Cyan Blue (Bullets/Icons)
BG_COLOR = RGBColor(255, 255, 255)      # Clean White (Content Background)
TEXT_COLOR_BODY = RGBColor(60, 60, 60)  # Dark Grey (Body Text)
TEXT_COLOR_COVER = RGBColor(255, 255, 255) # White (Cover Text)

# 2. Typography
FONT_HEAD = 'Arial'  # Headings (Bold)
FONT_BODY = 'Arial'  # Body (Regular)
FONT_SIZE_TITLE = Pt(24)
FONT_SIZE_BODY = Pt(11) # Spec: 10-12

# 3. Footer
FOOTER_TEXT = "Strictly Private & Confidential – Prepared by Kelp M&A Team"

def download_image(query, filename="temp_img.jpg"):
    """Downloads a generic stock image with browser headers."""
    if not query: return None
    safe_query = query.replace(" ", ",").lower()
    url = f"https://loremflickr.com/800/600/{safe_query}"
    print(f"   [Img] Downloading: {url}...")
    
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }
    
    try:
        response = requests.get(url, headers=headers, timeout=15)
        if response.status_code == 200 and len(response.content) > 1000:
            with open(filename, 'wb') as f:
                f.write(response.content)
            return filename
    except Exception as e:
        print(f"   [Warning] Image download error: {e}")
    return None

def create_cover_slide(prs, title_text="Investment Teaser"):
    """
    Creates the 'Dark Indigo' cover slide with geometric overlays.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # 1. Background Fill (Dark Indigo)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # 2. Geometric Overlay (Cyan Right Triangle)
    # Adds a subtle design element to the right side
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 
        Inches(6.5), Inches(0), Inches(3.5), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_COLOR
    shape.fill.transparency = 0.8  # High transparency
    shape.line.fill.background()   # No line

    # 3. Logo Placeholder (Top Right)
    logo_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.5), Inches(2), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "KELP"
    p.font.name = FONT_HEAD
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR_COVER
    p.alignment = PP_ALIGN.RIGHT

    # 4. Main Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Apex\n" + title_text
    p.font.name = FONT_HEAD
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = TEXT_COLOR_COVER
    
    # 5. Footer/Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(1))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = FOOTER_TEXT
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = ACCENT_COLOR 

def create_disclaimer_slide(prs):
    """
    Adds a legal disclaimer slide at the end.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Disclaimer"
    p.font.name = FONT_HEAD
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = PRIMARY_COLOR
    
    # Divider
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.2), Inches(9.5), Inches(1.2))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(2)

    # Disclaimer Text
    disclaimer_text = (
        "This presentation is strictly confidential and is being provided to you solely for your information. "
        "By accepting this presentation, you agree to keep it confidential and not to distribute it to any other person.\n\n"
        "This document does not constitute an offer to sell or a solicitation of an offer to buy any securities. "
        "The information contained herein has been obtained from sources believed to be reliable but is not guaranteed as to accuracy or completeness. "
        "Project Apex and Kelp M&A Team make no representation or warranty regarding the accuracy of the information.\n\n"
        "This presentation may contain forward-looking statements that involve risks and uncertainties. "
        "Actual results may differ materially from those projected."
    )
    
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    p = tf_body.paragraphs[0]
    p.text = disclaimer_text
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR_BODY
    p.alignment = PP_ALIGN.JUSTIFY

    # Add Footer
    apply_footer(slide)

def apply_footer(slide):
    """Adds the specific footer to bottom center."""
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(9), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.name = FONT_BODY
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER # Centered as requested

def apply_slide_branding(slide, slide_title):
    """
    Applies the Clean White background, Indigo Header, Pink/Orange Divider, and Logo.
    """
    # 1. Logo (Top Right) - Text Placeholder
    logo_box = slide.shapes.add_textbox(Inches(8), Inches(0.2), Inches(1.5), Inches(0.5))
    tf_logo = logo_box.text_frame
    p_logo = tf_logo.paragraphs[0]
    p_logo.text = "KELP"
    p_logo.font.name = FONT_HEAD
    p_logo.font.bold = True
    p_logo.font.size = Pt(14)
    p_logo.font.color.rgb = PRIMARY_COLOR
    p_logo.alignment = PP_ALIGN.RIGHT

    # 2. Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(7.5), Inches(0.8))
    title_tf = title_shape.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = slide_title.upper() 
    title_p.font.name = FONT_HEAD
    title_p.font.bold = True
    title_p.font.size = FONT_SIZE_TITLE
    title_p.font.color.rgb = PRIMARY_COLOR
    
    # 3. Divider Line
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 
        Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1)
    )
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)

    # 4. Footer
    apply_footer(slide)

def create_native_chart(slide, chart_data_dict):
    if not chart_data_dict or not chart_data_dict.get("values"): return 
    chart_data = CategoryChartData()
    chart_data.categories = chart_data_dict.get("labels", [])
    chart_data.add_series('Series 1', chart_data_dict.get("values", []))
    
    # Position: Right Quadrant
    x, y, cx, cy = Inches(5.5), Inches(1.5), Inches(4.0), Inches(3.5)
    try:
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
        chart.chart_title.text_frame.text = chart_data_dict.get("title", "Financial Overview")
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_HEAD
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    except: pass

def generate_styled_ppt(ppt_data, output_file):
    prs = Presentation()
    
    # 1. CREATE COVER SLIDE
    create_cover_slide(prs, "Investment Opportunity")
    
    slides_list = ppt_data.get("slides", [])
    for i, slide_content in enumerate(slides_list):
        # 2. CREATE CONTENT SLIDE
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        
        apply_slide_branding(slide, slide_content.get("title", "Slide"))
        
        # ---------------------------------------------------------
        # LEFT QUADRANT: Text
        # ---------------------------------------------------------
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.8), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        bullets = slide_content.get("bullets", [])
        for index, bullet in enumerate(bullets):
            p = tf.add_paragraph()
            # Anonymity Filter
            raw_text = bullet.get("text") or bullet.get("textiname") or bullet.get("summary")
            
            p.text = "■ " + raw_text 
            
            p.font.name = FONT_BODY
            p.font.size = FONT_SIZE_BODY
            p.font.color.rgb = TEXT_COLOR_BODY
            p.space_after = Pt(12) 

        # ---------------------------------------------------------
        # RIGHT QUADRANT: Visuals (Full Bleed-ish)
        # ---------------------------------------------------------
        chart_data = slide_content.get("chart_data")
        has_chart = chart_data and chart_data.get("values") and len(chart_data.get("values")) > 0

        if has_chart:
             create_native_chart(slide, chart_data)
        elif slide_content.get("image_query"):
             img_filename = f"temp_img_{i}.jpg"
             img_path = download_image(slide_content["image_query"], img_filename)
             if img_path:
                 try:
                    # 'Full Bleed' effect on the right edge
                    slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.5), width=Inches(4.5))
                 except Exception as e:
                    print(f"   [Error] Add picture failed: {e}")

    # 3. ADD DISCLAIMER SLIDE
    create_disclaimer_slide(prs)

    prs.save(output_file)
    print(f"PPT Saved: {output_file}")
    
    # Cleanup
    for i in range(len(slides_list)):
        if os.path.exists(f"temp_img_{i}.jpg"):
            try: os.remove(f"temp_img_{i}.jpg")
            except: pass