from llms import model
import json
from pptx import Presentation

def read_company_data(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def generate_ppt(ppt_points):
    prs = Presentation()
    data = ppt_points  # already a dict

    for slide_data in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = slide_data["title"]
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for bullet in slide_data["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 1

    prs.save("Blind_Investment_Teaser.pptx")



if __name__ == "__main__":
    data = read_company_data(
        "Company Data/automotive-kalyani-forge/Kalyani Forge-OnePager.md"
    )
    print("Data loaded")
    structured_output = model.get_response_from_llm(
        "llms/prompts/analyzer.txt",
        data
    )
    print(structured_output)

    ppt_points = model.get_response_from_llm(
        "llms/prompts/slide_gen.txt",
        structured_output
    )

    generate_ppt(ppt_points)

