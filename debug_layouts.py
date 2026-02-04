from pptx import Presentation

# CONFIG: Your Template Filename
TEMPLATE = "Sample Output.pptx" 

prs = Presentation(TEMPLATE)

print(f"--- ANALYZING {TEMPLATE} ---")
for i, layout in enumerate(prs.slide_layouts):
    print(f"Index {i}: {layout.name}")
    
    # Create a dummy slide for every layout
    slide = prs.slides.add_slide(layout)
    try:
        slide.shapes.title.text = f"Layout Index {i}"
    except:
        pass
        
    # Add a subtitle or body text if possible to see where it lands
    try:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = f"This is Body Text for Layout {i}"
    except:
        pass

prs.save("Layout_Test.pptx")
print("✅ Generated 'Layout_Test.pptx'. Open it and find the best slide style.")